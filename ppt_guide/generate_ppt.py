"""
WON AI 뱅커 PPT — 우리은행 가이드 톤 정밀 복제 버전
ppt_guide/*.jpg 디자인 시그니처를 카피:
 - 옅은 파랑 그라데이션 배경
 - 좌상단 우리은행 로고 + 좌하단 워터마크
 - 좌측 큰 헤드라인 ("우리" 등 핵심만 파란색) + 헤드라인 아래 짧은 파란 강조선
 - 우측 카드: 흰 배경 · 라운드 · 옅은 파란 아이콘 박스 · 점선 구분 · 체크리스트
 - 우측 1x4 카드: 좌측 원형 아이콘 + 번호(01~04) + 제목 + > 화살표
 - 좌하단 캐릭터 placeholder + 도시 실루엣
 - 우하단 막대그래프 장식
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── 컬러 ────────────────────────────────────────────────
WOORI_BLUE      = RGBColor(0x1b, 0x64, 0xda)
WOORI_BLUE_DARK = RGBColor(0x12, 0x49, 0xa3)
WOORI_BLUE_300  = RGBColor(0x60, 0xa5, 0xfa)
WOORI_BLUE_LITE = RGBColor(0xdb, 0xea, 0xfe)
BG_TOP          = RGBColor(0xff, 0xff, 0xff)
BG_BOT          = RGBColor(0xe8, 0xf2, 0xff)
GREEN           = RGBColor(0x22, 0xc5, 0x5e)
GREEN_LITE      = RGBColor(0xdc, 0xfc, 0xe7)
PURPLE          = RGBColor(0x8b, 0x5c, 0xf6)
PURPLE_LITE     = RGBColor(0xed, 0xe9, 0xfe)
ORANGE          = RGBColor(0xf5, 0x9e, 0x0b)
ORANGE_LITE     = RGBColor(0xfe, 0xf3, 0xc7)
RED             = RGBColor(0xdc, 0x26, 0x26)
RED_LITE        = RGBColor(0xfe, 0xe2, 0xe2)
GRAY_900        = RGBColor(0x11, 0x18, 0x27)
GRAY_700        = RGBColor(0x37, 0x41, 0x51)
GRAY_500        = RGBColor(0x6b, 0x72, 0x80)
GRAY_400        = RGBColor(0x9c, 0xa3, 0xaf)
GRAY_300        = RGBColor(0xd1, 0xd5, 0xdb)
GRAY_200        = RGBColor(0xe5, 0xe7, 0xeb)
GRAY_100        = RGBColor(0xf3, 0xf4, 0xf6)
GRAY_50         = RGBColor(0xf9, 0xfa, 0xfb)
WHITE           = RGBColor(0xff, 0xff, 0xff)
NAVY            = RGBColor(0x0f, 0x1b, 0x33)

FONT_KR = '맑은 고딕'

# ── 프레젠테이션 ────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
blank = prs.slide_layouts[6]


# ════════════════════════════════════════════════════════
# 공통 헬퍼
# ════════════════════════════════════════════════════════
def add_rect(s, l, t, w, h, fill=None, line=None, lw=None, shadow=False):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        if lw: sh.line.width = lw
    return sh


def add_round(s, l, t, w, h, fill=None, line=None, lw=None, radius=0.06):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.adjustments[0] = radius
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        if lw: sh.line.width = lw
    return sh


def add_oval(s, l, t, w, h, fill, line=None, lw=None):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    sh.shadow.inherit = False
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        if lw: sh.line.width = lw
    return sh


def add_text(s, l, t, w, h, text, size=14, bold=False, color=GRAY_700,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_sp=1.3, font=FONT_KR):
    box = s.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_sp
        run = p.add_run(); run.text = line
        run.font.name = font; run.font.size = Pt(size)
        run.font.bold = bold; run.font.color.rgb = color
    return box


def add_rich(s, l, t, w, h, segs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_sp=1.3):
    """segs: [(text, size, bold, color), ...] or 'NL'"""
    box = s.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align; p.line_spacing = line_sp
    for seg in segs:
        if seg == 'NL':
            p = tf.add_paragraph(); p.alignment = align; p.line_spacing = line_sp
            continue
        text, size, bold, color = seg
        run = p.add_run(); run.text = text
        run.font.name = FONT_KR; run.font.size = Pt(size)
        run.font.bold = bold; run.font.color.rgb = color
    return box


def add_line(s, l, t, w, h, color=GRAY_300, lw=Pt(1)):
    """직선 (얇은 사각형)"""
    return add_rect(s, l, t, w, h, fill=color)


def add_dotted_line(s, l, t, w, color=GRAY_300):
    """가로 점선 (작은 동그라미 반복)"""
    dot_size = Emu(20000)
    gap = Emu(50000)
    x = l
    while x < l + w:
        add_oval(s, x, t, dot_size, dot_size, color)
        x += dot_size + gap


def add_gradient_bg(s):
    """옅은 파랑 그라데이션 (상단 흰색 → 하단 파랑) — XML 조작"""
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.shadow.inherit = False
    bg.line.fill.background()
    spPr = bg.fill._xPr
    # 기존 fill 노드 제거
    for tag in ('a:solidFill', 'a:gradFill', 'a:noFill', 'a:pattFill', 'a:blipFill'):
        for child in spPr.findall(qn(tag)):
            spPr.remove(child)
    # gradFill 직접 추가
    gradFill_xml = (
        '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
        'flip="none" rotWithShape="1">'
        '<a:gsLst>'
        '<a:gs pos="0"><a:srgbClr val="FFFFFF"/></a:gs>'
        '<a:gs pos="100000"><a:srgbClr val="E8F2FF"/></a:gs>'
        '</a:gsLst>'
        '<a:lin ang="5400000" scaled="1"/>'
        '</a:gradFill>'
    )
    grad_elem = etree.fromstring(gradFill_xml)
    # line 앞에 삽입
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        ln.addprevious(grad_elem)
    else:
        spPr.append(grad_elem)
    return bg


def add_solid_bg(s, color=WHITE):
    return add_rect(s, 0, 0, SW, SH, fill=color)


def add_header_logo(s):
    """좌상단 우리은행 로고"""
    # 파란 원형 마크
    add_oval(s, Inches(0.55), Inches(0.42), Inches(0.32), Inches(0.32), WOORI_BLUE)
    add_oval(s, Inches(0.61), Inches(0.48), Inches(0.20), Inches(0.20), WHITE)
    # "우리은행" 텍스트
    add_text(s, Inches(0.95), Inches(0.4), Inches(2.0), Inches(0.4),
        "우리은행", size=18, bold=True, color=GRAY_900,
        anchor=MSO_ANCHOR.MIDDLE)


def add_page_no(s, num):
    """우상단 페이지 번호"""
    add_text(s, Inches(11.7), Inches(0.42), Inches(0.6), Inches(0.4),
        f"{num:02d}", size=14, bold=True, color=WOORI_BLUE,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    add_line(s, Inches(12.4), Inches(0.60), Inches(0.45), Emu(15000), color=WOORI_BLUE)


def add_footer(s, label="2026.05  |  WOORI BANK"):
    """좌하단 워터마크"""
    add_text(s, Inches(0.55), Inches(7.05), Inches(4), Inches(0.4),
        label, size=11, color=GRAY_400, anchor=MSO_ANCHOR.MIDDLE)


def add_bar_chart_deco(s):
    """우하단 작은 막대그래프 장식"""
    base_l, base_t = Inches(12.55), Inches(6.95)
    heights = [Inches(0.18), Inches(0.27), Inches(0.36)]
    for i, h in enumerate(heights):
        add_round(s, base_l + Inches(0.18 * i), base_t + Inches(0.4) - h,
            Inches(0.13), h, fill=WOORI_BLUE_300, radius=0.15)


def add_corner_decor(s):
    """좌하단 옅은 곡선 장식 (구름 같은)"""
    # 옅은 파랑 반원
    add_oval(s, Inches(-1.5), Inches(6.5), Inches(4), Inches(2.5),
        fill=WOORI_BLUE_LITE)


def add_underline(s, l, t, w=Inches(0.6), color=WOORI_BLUE, h=Emu(40000)):
    """헤드라인 아래 짧은 강조선"""
    return add_rect(s, l, t, w, h, fill=color)


def add_check_item(s, l, t, w, text, size=12, color=GRAY_700):
    """체크리스트 아이템: 파란 원 + 흰 체크 + 텍스트"""
    add_oval(s, l, t + Inches(0.05), Inches(0.2), Inches(0.2), WOORI_BLUE)
    add_text(s, l, t + Inches(0.05), Inches(0.2), Inches(0.2),
        "✓", size=10, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, l + Inches(0.32), t, w - Inches(0.32), Inches(0.3),
        text, size=size, color=color, anchor=MSO_ANCHOR.MIDDLE)


def add_icon_box(s, l, t, size, bg_color, icon_color, icon_text="🛡"):
    """둥근 사각형 아이콘 박스 (옅은 배경 + 진한 아이콘)"""
    add_round(s, l, t, size, size, fill=bg_color, radius=0.2)
    add_text(s, l, t, size, size, icon_text, size=int(size.inches * 24),
        color=icon_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_card(s, l, t, w, h, with_shadow=True):
    """기본 카드 (흰 배경 + 라운드 + 그림자)"""
    card = add_round(s, l, t, w, h, fill=WHITE, radius=0.04)
    if with_shadow:
        # 그림자는 PPT shadow effect로 (간단히 회색 사각형 뒤에 깔기로 대체)
        pass
    return card


def add_chapter_label(s, l, t, text, color=WOORI_BLUE):
    """슬라이드 챕터 라벨 (헤드라인 위 작은 글씨)"""
    add_text(s, l, t, Inches(4), Inches(0.3),
        f"●  {text}", size=11, bold=True, color=color)


def add_character_placeholder(s, l, t, w, h):
    """좌하단 캐릭터 자리 (가이드 jpg에서 잘라 넣을 영역)"""
    # 옅은 점선 박스
    add_round(s, l, t, w, h, fill=WOORI_BLUE_LITE, radius=0.08)
    add_text(s, l, t, w, h - Inches(0.3),
        "🐻  🐰  🐝  🐻", size=42, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, l, t + h - Inches(0.4), w, Inches(0.3),
        "[ 우리은행 캐릭터 자리 ]", size=9, color=WOORI_BLUE_300,
        align=PP_ALIGN.CENTER)


def add_city_silhouette(s, l, t, w, h):
    """좌하단 도시 실루엣 (옅은 파랑 사각형들)"""
    blocks = [
        (0.0, 0.6, 0.15, 0.4),
        (0.13, 0.4, 0.15, 0.6),
        (0.26, 0.7, 0.10, 0.3),
        (0.34, 0.3, 0.12, 0.7),
        (0.44, 0.55, 0.18, 0.45),
        (0.60, 0.45, 0.10, 0.55),
        (0.68, 0.65, 0.13, 0.35),
        (0.79, 0.5, 0.12, 0.5),
        (0.89, 0.7, 0.11, 0.3),
    ]
    for (x_pct, y_pct, w_pct, h_pct) in blocks:
        add_rect(s, l + int(w * x_pct), t + int(h * y_pct),
            int(w * w_pct), int(h * h_pct),
            fill=WOORI_BLUE_LITE)


# ════════════════════════════════════════════════════════
# 슬라이드 1 — 표지
# ════════════════════════════════════════════════════════
def slide_01():
    s = prs.slides.add_slide(blank)
    add_gradient_bg(s)
    add_header_logo(s)

    # 좌측 상단 작은 카피
    add_text(s, Inches(0.7), Inches(1.6), Inches(6), Inches(0.35),
        "옆자리 행원이 사라진 시대,  ", size=14, color=GRAY_500)
    add_text(s, Inches(0.7), Inches(1.6), Inches(6), Inches(0.35),
        "                                      불완전판매 예방 AI", size=14, color=WOORI_BLUE, bold=True)

    # 메인 헤드라인 — 가이드 톤 (매우 굵게, "우리"만 파랑)
    add_rich(s, Inches(0.7), Inches(2.1), Inches(8), Inches(2.5),
        segs=[
            ('옆에 행원이 없으면,', 54, True, GRAY_900),
            'NL',
            ('우리', 54, True, WOORI_BLUE),
            ('의 다음 클릭이', 54, True, GRAY_900),
            'NL',
            ('위험합니다.', 54, True, GRAY_900),
        ], line_sp=1.15)

    # 강조선
    add_underline(s, Inches(0.7), Inches(5.0), w=Inches(0.6))

    # 부연 설명
    add_text(s, Inches(0.7), Inches(5.2), Inches(7), Inches(1.5),
        "WON AI 뱅커는 마이데이터 기반 개인화 경고와\n음성 안내로, 사라진 옆자리 행원의 역할을\n디지털에 다시 앉힙니다.",
        size=15, color=GRAY_500, line_sp=1.7)

    # 우측 폰 화면 미리보기 (3개 겹쳐서)
    phone_w, phone_h = Inches(2.5), Inches(4.8)
    # 뒷 폰 1
    add_round(s, Inches(8.0), Inches(1.4), phone_w, phone_h,
        fill=WHITE, line=GRAY_200, lw=Pt(0.5), radius=0.05)
    add_text(s, Inches(8.0), Inches(1.5), phone_w, Inches(0.3),
        "투자성향 분석", size=10, bold=True, color=GRAY_700, align=PP_ALIGN.CENTER)
    for i, q in enumerate(['Q. 손실 감수도?', 'Q. 투자 경험은?', 'Q. 투자 목표?']):
        add_round(s, Inches(8.15), Inches(2.0 + 0.5 * i), phone_w - Inches(0.3), Inches(0.4),
            fill=GRAY_50, radius=0.15)
        add_text(s, Inches(8.3), Inches(2.05 + 0.5 * i), phone_w - Inches(0.5), Inches(0.3),
            q, size=9, color=GRAY_700, anchor=MSO_ANCHOR.MIDDLE)

    # 가운데 폰 2
    add_round(s, Inches(9.6), Inches(1.2), phone_w, phone_h,
        fill=WHITE, line=GRAY_200, lw=Pt(0.5), radius=0.05)
    add_text(s, Inches(9.6), Inches(1.3), phone_w, Inches(0.3),
        "AI 사전 경고", size=10, bold=True, color=WOORI_BLUE, align=PP_ALIGN.CENTER)
    # 챗 메시지
    add_round(s, Inches(9.75), Inches(1.85), phone_w - Inches(0.3), Inches(1.8),
        fill=GRAY_50, radius=0.08)
    add_text(s, Inches(9.9), Inches(2.0), phone_w - Inches(0.6), Inches(1.6),
        "⚠ 경고 1. 유동성\n\n이혜원님,\n자동이체 1,405,000원\n예정. 30만원 묶으면\n잔액 부족 위험.",
        size=8.5, color=GRAY_700, line_sp=1.45)
    # 버튼
    add_round(s, Inches(9.75), Inches(3.95), phone_w - Inches(0.3), Inches(0.45),
        fill=WHITE, line=RED, lw=Pt(1.5), radius=0.2)
    add_text(s, Inches(9.75), Inches(3.95), phone_w - Inches(0.3), Inches(0.45),
        "다시 생각해볼게요", size=10, bold=True, color=RED,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_round(s, Inches(9.75), Inches(4.5), phone_w - Inches(0.3), Inches(0.45),
        fill=GRAY_100, radius=0.2)
    add_text(s, Inches(9.75), Inches(4.5), phone_w - Inches(0.3), Inches(0.45),
        "네, 그래도 진행", size=10, color=GRAY_500,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 앞 폰 3
    add_round(s, Inches(11.0), Inches(1.0), phone_w, phone_h,
        fill=WHITE, line=GRAY_200, lw=Pt(0.5), radius=0.05)
    add_text(s, Inches(11.0), Inches(1.1), phone_w, Inches(0.3),
        "📊 마이데이터", size=10, bold=True, color=WOORI_BLUE, align=PP_ALIGN.CENTER)
    # 차트 막대
    add_text(s, Inches(11.15), Inches(1.5), phone_w - Inches(0.3), Inches(0.3),
        "이번 달 월급 흐름", size=8, color=GRAY_500)
    bars = [(0.42, GRAY_500), (0.20, ORANGE), (0.08, GREEN), (0.12, WOORI_BLUE)]
    bx = Inches(11.15)
    for w_ratio, color in bars:
        bw = int((phone_w - Inches(0.3)) * w_ratio)
        add_rect(s, bx, Inches(1.85), bw, Inches(0.22), fill=color)
        bx += bw
    add_text(s, Inches(11.15), Inches(2.15), phone_w - Inches(0.3), Inches(0.3),
        "고정 변동 적금 여유", size=7, color=GRAY_500)
    # 가입 직전 점검
    add_round(s, Inches(11.15), Inches(2.7), phone_w - Inches(0.3), Inches(1.9),
        fill=RED_LITE, line=RED, lw=Pt(1), radius=0.06)
    add_text(s, Inches(11.25), Inches(2.8), phone_w - Inches(0.5), Inches(0.3),
        "⚠ 꼭 확인하세요", size=10, bold=True, color=RED)
    add_text(s, Inches(11.25), Inches(3.1), phone_w - Inches(0.5), Inches(1.4),
        "손실 시나리오\n\n300,000원\n→ 198,000원\n-34%",
        size=9, color=GRAY_700, line_sp=1.5)

    # 좌하단 캐릭터 + 도시 (가이드 톤)
    add_corner_decor(s)

    # 하단 라벨
    add_footer(s, "우리은행 × SSAFY AI-금융소비자보호 아이디어 경진대회 2026  ·  팀 [팀명]")
    add_bar_chart_deco(s)


# ════════════════════════════════════════════════════════
# 슬라이드 2 — Scene: 사라진 옆자리
# ════════════════════════════════════════════════════════
def slide_02():
    s = prs.slides.add_slide(blank)
    add_gradient_bg(s)
    add_header_logo(s)
    add_page_no(s, 2)
    add_footer(s)
    add_bar_chart_deco(s)

    add_chapter_label(s, Inches(0.7), Inches(1.4), "SCENE  01")

    add_rich(s, Inches(0.7), Inches(1.9), Inches(8), Inches(2.5),
        segs=[
            ('5년 사이,', 50, True, GRAY_900),
            'NL',
            ('동네 은행이 ', 50, True, GRAY_900),
            ('사라졌습니다.', 50, True, WOORI_BLUE),
        ], line_sp=1.2)

    add_underline(s, Inches(0.7), Inches(4.7))

    add_text(s, Inches(0.7), Inches(4.95), Inches(7), Inches(1.2),
        "노년층 요구는 단순 잔액조회가 아닙니다.\n연금 수령, 병원비 이체, 정기지출 설정 같은\n일상 밀착형 금융 — 점포 사라지면 곧 금융 접근권 차단.",
        size=14, color=GRAY_500, line_sp=1.7)

    add_text(s, Inches(0.7), Inches(6.6), Inches(8), Inches(0.3),
        "출처. 한국경제 「은행 점포 폐쇄 가속, 작년 100개 넘게 줄었다」 (2026.01)",
        size=10, color=GRAY_400)

    # 우측 큰 숫자 카드
    box_l, box_t = Inches(8.5), Inches(1.9)
    box_w, box_h = Inches(4.3), Inches(5.0)
    add_round(s, box_l, box_t, box_w, box_h,
        fill=WHITE, line=GRAY_200, lw=Pt(0.5), radius=0.04)

    add_text(s, box_l, box_t + Inches(0.4), box_w, Inches(0.4),
        "최근 5년 은행 점포 변화", size=13, bold=True, color=GRAY_700,
        align=PP_ALIGN.CENTER)

    add_text(s, box_l, box_t + Inches(0.9), box_w, Inches(1.5),
        "904", size=120, bold=True, color=WOORI_BLUE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, box_l, box_t + Inches(2.5), box_w, Inches(0.4),
        "개 폐쇄  ·  -14.1%", size=18, bold=True, color=GRAY_900,
        align=PP_ALIGN.CENTER)

    add_dotted_line(s, box_l + Inches(0.3), box_t + Inches(3.15), box_w - Inches(0.6))

    # 막대 비교
    add_text(s, box_l + Inches(0.3), box_t + Inches(3.35), Inches(1.8), Inches(0.3),
        "2020년 말", size=10, color=GRAY_500)
    add_text(s, box_l + Inches(0.3), box_t + Inches(3.6), Inches(2), Inches(0.4),
        "6,427개", size=18, bold=True, color=GRAY_700)
    add_rect(s, box_l + Inches(2.4), box_t + Inches(3.55),
        Inches(1.6), Inches(0.45), fill=WOORI_BLUE_LITE)

    add_text(s, box_l + Inches(0.3), box_t + Inches(4.1), Inches(1.8), Inches(0.3),
        "2025년 9월", size=10, color=GRAY_500)
    add_text(s, box_l + Inches(0.3), box_t + Inches(4.35), Inches(2), Inches(0.4),
        "5,523개", size=18, bold=True, color=RED)
    add_rect(s, box_l + Inches(2.4), box_t + Inches(4.3),
        Inches(1.37), Inches(0.45), fill=RED)


# ════════════════════════════════════════════════════════
# 슬라이드 3 — 70대 디지털 절벽 (좌헤드 + 우 2x2 카드)
# ════════════════════════════════════════════════════════
def slide_03():
    s = prs.slides.add_slide(blank)
    add_gradient_bg(s)
    add_header_logo(s)
    add_page_no(s, 3)
    add_footer(s)
    add_bar_chart_deco(s)

    add_chapter_label(s, Inches(0.7), Inches(1.4), "PROBLEM  01  ·  디지털 격차")

    add_rich(s, Inches(0.7), Inches(1.9), Inches(7), Inches(2),
        segs=[
            ('60대는 카톡을 합니다.', 38, True, GRAY_900),
            'NL',
            ('70', 38, True, WOORI_BLUE),
            ('대는 ', 38, True, GRAY_900),
            ('인터넷', 38, True, RED),
            ('을 못 합니다.', 38, True, GRAY_900),
        ], line_sp=1.25)

    add_underline(s, Inches(0.7), Inches(4.2))

    add_text(s, Inches(0.7), Inches(4.4), Inches(6.5), Inches(2.3),
        "같은 \"고령층\" 안에 거대한 절벽이 있습니다.\nKISDI 2024 조사에서 70대의 인터넷 활용은\n5점 만점에 1.85점, 거의 낙제 수준입니다.\n\n그런데도 65세 이상 38%는\n\"은행앱 쓰고 싶다\"고 답했습니다.",
        size=14, color=GRAY_500, line_sp=1.8)

    # 우측: 우리가 분석한 격차 (가이드2 카드 형식)
    add_text(s, Inches(7.5), Inches(1.4), Inches(5.5), Inches(0.4),
        "우리가 본 디지털 절벽", size=14, bold=True, color=GRAY_900)
    add_line(s, Inches(10.5), Inches(1.65), Inches(2.4), Emu(8000), color=GRAY_200)

    # 2x2 카드
    cards = [
        ('📱', '메신저 격차',  WOORI_BLUE_LITE, WOORI_BLUE,
         '60대 4.23점 — 국민 평균 수준\n70대 1점대 — 낙제',
         ['60대는 카톡·문자 OK', '70대는 통화 위주', 'KISDI 2024']),
        ('🌐', '인터넷 격차',  RED_LITE, RED,
         '60대 3.39점 → 70대 1.85점\n격차 1.54점 (절벽)',
         ['단순 검색도 어려움', '뉴스·정보 차단', '온라인 일상 불가']),
        ('💰', '금융앱 격차',  ORANGE_LITE, ORANGE,
         '실제 모바일뱅킹 이용\n60대 약 30% / 70대 한자릿수',
         ['가입 ≠ 실사용', '진입 단계서 포기', '갤럽 2025']),
        ('💬', '욕구는 있다', GREEN_LITE, GREEN,
         '65세 이상 38%가\n"은행앱 쓰고 싶다"',
         ['수요는 분명', '진입 못 할 뿐', '국제신문 2025']),
    ]
    base_l = Inches(7.5)
    base_t = Inches(1.9)
    card_w = Inches(2.65)
    card_h = Inches(2.45)
    gap = Inches(0.15)
    for i, (icon, title, ibg, ic, desc, checks) in enumerate(cards):
        row, col = i // 2, i % 2
        x = base_l + (card_w + gap) * col
        y = base_t + (card_h + gap) * row
        add_round(s, x, y, card_w, card_h, fill=WHITE, line=GRAY_200, lw=Pt(0.5), radius=0.04)
        add_icon_box(s, x + Inches(0.2), y + Inches(0.2), Inches(0.55), ibg, ic, icon)
        add_text(s, x + Inches(0.2), y + Inches(0.88), card_w - Inches(0.4), Inches(0.35),
            title, size=14, bold=True, color=WOORI_BLUE)
        add_text(s, x + Inches(0.2), y + Inches(1.25), card_w - Inches(0.4), Inches(0.6),
            desc, size=9.5, color=GRAY_500, line_sp=1.4)
        add_dotted_line(s, x + Inches(0.2), y + Inches(1.85), card_w - Inches(0.4))
        for j, ck in enumerate(checks):
            add_check_item(s, x + Inches(0.2), y + Inches(2.0) + Inches(0.15 * j),
                card_w - Inches(0.4), ck, size=8.5, color=GRAY_500)


# ════════════════════════════════════════════════════════
# 슬라이드 4 — 4단계 진입 장벽 (1x4 가이드3 패턴)
# ════════════════════════════════════════════════════════
def slide_04():
    s = prs.slides.add_slide(blank)
    add_gradient_bg(s)
    add_header_logo(s)
    add_page_no(s, 4)
    add_footer(s)
    add_bar_chart_deco(s)

    add_chapter_label(s, Inches(0.7), Inches(1.4), "PROBLEM  02  ·  진입 장벽")

    add_rich(s, Inches(0.7), Inches(1.9), Inches(6.5), Inches(2.5),
        segs=[
            ('네 단계 ', 44, True, GRAY_900),
            ('모두', 44, True, RED),
            ('에서', 44, True, GRAY_900),
            'NL',
            ('막힙니다.', 44, True, GRAY_900),
        ], line_sp=1.2)

    add_underline(s, Inches(0.7), Inches(4.3))

    add_text(s, Inches(0.7), Inches(4.5), Inches(6), Inches(1.5),
        "쓸 줄 모르는 게 문제가 아닙니다.\n진입 단계마다 함정이 있고, 그 함정을 통과한\n사람만 \"쉬운 모드\"를 만날 수 있습니다.",
        size=14, color=GRAY_500, line_sp=1.7)

    # 좌하단 말풍선 박스
    bubble_l, bubble_t = Inches(0.7), Inches(6.0)
    add_round(s, bubble_l, bubble_t, Inches(6.0), Inches(0.85),
        fill=WHITE, line=WOORI_BLUE_LITE, lw=Pt(1), radius=0.1)
    add_oval(s, bubble_l + Inches(0.2), bubble_t + Inches(0.2),
        Inches(0.45), Inches(0.45), WOORI_BLUE)
    add_text(s, bubble_l + Inches(0.2), bubble_t + Inches(0.2),
        Inches(0.45), Inches(0.45), "💡", size=14, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, bubble_l + Inches(0.8), bubble_t + Inches(0.1),
        Inches(5), Inches(0.35),
        "결정적 한계", size=12, bold=True, color=WOORI_BLUE)
    add_text(s, bubble_l + Inches(0.8), bubble_t + Inches(0.4),
        Inches(5), Inches(0.4),
        "\"쉬운 모드도 도달해야 쓸 수 있다.\"",
        size=11, bold=True, color=GRAY_700)

    # 우측 1x4 카드
    add_text(s, Inches(7.4), Inches(1.4), Inches(5.5), Inches(0.4),
        "●  단계별로 어디서 막히는가", size=14, bold=True, color=GRAY_900)
    add_line(s, Inches(10.4), Inches(1.65), Inches(2.5), Emu(8000), color=GRAY_200)

    stages = [
        ('01', '진입 자체', WOORI_BLUE_LITE, WOORI_BLUE,
         '"쉬운 모드"를 켜려면 먼저 앱 설치 + 본인인증 + 로그인을 통과해야 합니다.'),
        ('02', '인증 함정', PURPLE_LITE, PURPLE,
         'OTP는 공동인증서로, 공동인증서는 OTP로 풀어야 하는 무한루프. 셀카 인증 실패도 빈번.'),
        ('03', '화면 오인지', ORANGE_LITE, ORANGE,
         '화살표 클릭 가능 인지 못 함 · 스크롤 모름 · 예시를 본인 데이터로 오해 (Toss UX 2024).'),
        ('04', '사기 노출', RED_LITE, RED,
         '진입은 못해도 송금은 함. 60대 보이스피싱 4.6배 폭증, 펀드 불완전판매 민원 30%.'),
    ]
    base_t = Inches(1.95)
    card_h = Inches(1.18)
    card_l = Inches(7.4)
    card_w = Inches(5.5)
    for i, (num, title, ibg, ic, desc) in enumerate(stages):
        y = base_t + (card_h + Inches(0.08)) * i
        add_round(s, card_l, y, card_w, card_h, fill=WHITE, line=GRAY_200, lw=Pt(0.5), radius=0.05)
        # 좌측 원형 아이콘
        add_oval(s, card_l + Inches(0.2), y + Inches(0.2), Inches(0.78), Inches(0.78), ibg)
        add_text(s, card_l + Inches(0.2), y + Inches(0.2), Inches(0.78), Inches(0.78),
            num, size=20, bold=True, color=ic,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 번호 (작게 위)
        add_text(s, card_l + Inches(1.2), y + Inches(0.15), Inches(0.5), Inches(0.25),
            num, size=10, bold=True, color=ic)
        # 제목
        add_text(s, card_l + Inches(1.2), y + Inches(0.35), card_w - Inches(1.5), Inches(0.35),
            title, size=15, bold=True, color=GRAY_900)
        # 설명
        add_text(s, card_l + Inches(1.2), y + Inches(0.72), card_w - Inches(1.5), Inches(0.5),
            desc, size=10, color=GRAY_500, line_sp=1.4)
        # 우측 > 화살표
        add_text(s, card_l + card_w - Inches(0.45), y, Inches(0.35), card_h,
            "›", size=22, color=ic, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════
# 슬라이드 5 — 사기 + 불완전판매 (3 박스 + 인용)
# ════════════════════════════════════════════════════════
def slide_05():
    s = prs.slides.add_slide(blank)
    add_gradient_bg(s)
    add_header_logo(s)
    add_page_no(s, 5)
    add_footer(s)
    add_bar_chart_deco(s)

    add_chapter_label(s, Inches(0.7), Inches(1.4), "PROBLEM  03  ·  피해의 무게")

    add_rich(s, Inches(0.7), Inches(1.9), Inches(11), Inches(1.2),
        segs=[
            ('숫자가 말해주는 ', 40, True, GRAY_900),
            ('피해 구조', 40, True, WOORI_BLUE),
        ])

    add_underline(s, Inches(0.7), Inches(3.0))

    add_text(s, Inches(0.7), Inches(3.2), Inches(12), Inches(0.5),
        "시니어·취약계층이 집중 타격을 받습니다.",
        size=16, color=GRAY_500)

    # 3개 큰 숫자 박스
    boxes = [
        ('30', '%', '펀드 불완전판매 민원\n60세 이상 비중',
         '금감원 (2020~2022)', WOORI_BLUE,
         ['은행 31.7% / 증권사 32.1%', '인구 비중보다 훨씬 높음', '금소법 사각지대']),
        ('4.6', '배', '60대 보이스피싱 피해 폭증\n2016년 1,261건 → 2025년 5,801건',
         '경찰청', RED,
         ['2025 상반기 비중 30.6%', '2020년의 약 2배', '시니어 집중 타깃화']),
        ('8,856', '억원', '2025년 8월 누적 피해액\n역대 최고치',
         '금융감독원', PURPLE,
         ['2023년 4,472억 → 2025 폭증', 'AI 딥페이크 사기 영향', '평균 피해도 증가']),
    ]
    box_t = Inches(3.9)
    box_w = Inches(4.0)
    box_h = Inches(2.5)
    gap = Inches(0.2)
    total_w = box_w * 3 + gap * 2
    start_l = (SW - total_w) // 2
    for i, (num, unit, label, src, color, checks) in enumerate(boxes):
        l = start_l + (box_w + gap) * i
        add_round(s, l, box_t, box_w, box_h, fill=WHITE, line=GRAY_200, lw=Pt(0.5), radius=0.04)
        # 큰 숫자 + 단위
        add_rich(s, l + Inches(0.3), box_t + Inches(0.3), box_w - Inches(0.6), Inches(1.0),
            segs=[
                (num, 56, True, color),
                (unit, 22, True, color),
            ], align=PP_ALIGN.CENTER)
        # 라벨
        add_text(s, l + Inches(0.3), box_t + Inches(1.3), box_w - Inches(0.6), Inches(0.6),
            label, size=11.5, color=GRAY_700, align=PP_ALIGN.CENTER, line_sp=1.4)
        # 점선
        add_dotted_line(s, l + Inches(0.3), box_t + Inches(1.95), box_w - Inches(0.6))
        # 체크리스트
        for j, ck in enumerate(checks):
            add_check_item(s, l + Inches(0.3), box_t + Inches(2.05) + Inches(0.14 * j),
                box_w - Inches(0.6), ck, size=8, color=GRAY_500)

    # 하단 빨간 인용 박스
    quote_t = Inches(6.55)
    add_round(s, Inches(0.7), quote_t, Inches(12), Inches(0.55),
        fill=RED_LITE, line=RED, lw=Pt(1), radius=0.1)
    add_rich(s, Inches(0.9), quote_t, Inches(11.7), Inches(0.55),
        segs=[
            ('⚠  ', 13, True, RED),
            ('AI 음성·딥페이크 진화  —  30초 음성만으로 가족 목소리 복제. 2025년 서울 직장인은 \'아들 얼굴\' 딥페이크 영상통화에 ', 11, False, GRAY_700),
            ('5,000만원 송금. ', 11, True, RED),
            ('시니어는 "딥페이크" 개념 자체를 모릅니다.', 11, False, GRAY_700),
        ], anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════
# 슬라이드 6 — 은행도 노력 중인데… (좌 노력 / 우 한계)
# ════════════════════════════════════════════════════════
def slide_06():
    s = prs.slides.add_slide(blank)
    add_gradient_bg(s)
    add_header_logo(s)
    add_page_no(s, 6)
    add_footer(s)
    add_bar_chart_deco(s)

    add_chapter_label(s, Inches(0.7), Inches(1.4), "PROBLEM  04  ·  기존 대응의 한계")

    add_rich(s, Inches(0.7), Inches(1.9), Inches(7), Inches(2),
        segs=[
            ('은행도 ', 40, True, GRAY_900),
            ('가만히', 40, True, WOORI_BLUE),
            ('있진 않았습니다.', 40, True, GRAY_900),
            'NL',
            ('그런데 ', 40, True, GRAY_900),
            ('충분하지 않습니다.', 40, True, RED),
        ], line_sp=1.25)

    add_underline(s, Inches(0.7), Inches(5.0))

    add_text(s, Inches(0.7), Inches(5.2), Inches(6.5), Inches(1.5),
        "가이드라인은 \"방향\"만 제시했고,\n실제 구현은 \"글자 크기 키우기\"에 머물렀습니다.\n\n결정적으로, 그 \"쉬운 모드\"는\n진입을 통과한 사람만 만날 수 있습니다.",
        size=13, color=GRAY_500, line_sp=1.7)

    # 우측: 2분할 카드
    # 위 — 은행 노력
    up_l, up_t = Inches(7.5), Inches(1.9)
    up_w, up_h = Inches(5.3), Inches(2.4)
    add_round(s, up_l, up_t, up_w, up_h, fill=WHITE, line=GRAY_200, lw=Pt(0.5), radius=0.04)
    add_text(s, up_l + Inches(0.3), up_t + Inches(0.25), up_w, Inches(0.4),
        "✅  은행권이 시도한 대응", size=14, bold=True, color=GRAY_900)
    efforts = [
        '4대 은행 모두 "고령자 모드" 출시 (2023)',
        '우리은행 시니어플러스 영업점 · 사랑채 운영',
        '카카오뱅크 · 토스 "간편홈" 도입',
        '금융위 "고령자 친화적 앱 구성지침" (2022)',
        '한도제한계좌 · 지급정지 · 안심차단 인프라',
    ]
    for i, txt in enumerate(efforts):
        add_check_item(s, up_l + Inches(0.4), up_t + Inches(0.75) + Inches(0.3 * i),
            up_w - Inches(0.6), txt, size=11, color=GRAY_700)

    # 아래 — 한계 (빨간)
    dn_l, dn_t = Inches(7.5), Inches(4.4)
    dn_w, dn_h = Inches(5.3), Inches(2.45)
    add_round(s, dn_l, dn_t, dn_w, dn_h, fill=RED_LITE, line=RED, lw=Pt(1.5), radius=0.04)
    add_text(s, dn_l + Inches(0.3), dn_t + Inches(0.25), dn_w, Inches(0.4),
        '⚠  그런데, 한 기사 제목이 모든 걸 말합니다', size=12, bold=True, color=RED)
    add_text(s, dn_l + Inches(0.3), dn_t + Inches(0.7), dn_w - Inches(0.6), Inches(0.8),
        '"고령자 모드 있다지만,\n쉬운 금융은 없었다."',
        size=18, bold=True, color=GRAY_900, line_sp=1.3)
    add_text(s, dn_l + Inches(0.3), dn_t + Inches(1.65), dn_w - Inches(0.6), Inches(0.4),
        '— 브라보마이라이프, 4대 은행 앱 조사 (2025)',
        size=10, color=GRAY_500)
    add_dotted_line(s, dn_l + Inches(0.3), dn_t + Inches(2.0), dn_w - Inches(0.6), color=RED)
    add_text(s, dn_l + Inches(0.3), dn_t + Inches(2.1), dn_w - Inches(0.6), Inches(0.4),
        '쉬운 모드도 "도달해야" 쓸 수 있다 → 진입 차단',
        size=11, bold=True, color=RED)


# ════════════════════════════════════════════════════════
# 슬라이드 7 — 그래서 (전환 — 솔루션 등장)
# ════════════════════════════════════════════════════════
def slide_07():
    s = prs.slides.add_slide(blank)
    add_gradient_bg(s)
    add_header_logo(s)
    add_page_no(s, 7)
    add_footer(s)
    add_bar_chart_deco(s)

    add_chapter_label(s, Inches(0.7), Inches(1.4), "SOLUTION  ·  옆자리 행원의 귀환")

    # 중앙 큰 헤드라인
    add_text(s, Inches(0.7), Inches(2.2), Inches(12), Inches(1.0),
        "그래서 저희는 ─", size=22, color=GRAY_500, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(2.8), Inches(12), Inches(1.5),
        "WON AI 뱅커", size=88, bold=True, color=WOORI_BLUE,
        align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(4.5), Inches(12), Inches(0.5),
        "사라진 옆자리에 다시 앉습니다.",
        size=22, color=GRAY_700, align=PP_ALIGN.CENTER)

    # 3줄 카피
    line_t = Inches(5.5)
    add_text(s, Inches(0.7), line_t, Inches(12), Inches(0.45),
        "권유하지 않습니다.", size=22, bold=True, color=GRAY_900, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), line_t + Inches(0.45), Inches(12), Inches(0.45),
        "같이 고민합니다.", size=22, bold=True, color=WOORI_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), line_t + Inches(0.9), Inches(12), Inches(0.45),
        "끝까지 함께합니다.", size=22, bold=True, color=GRAY_900, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# 슬라이드 8 — 4가지 핵심 가치 (가이드2 정확 카피)
# ════════════════════════════════════════════════════════
def slide_08():
    s = prs.slides.add_slide(blank)
    add_gradient_bg(s)
    add_header_logo(s)
    add_page_no(s, 8)
    add_footer(s)
    add_bar_chart_deco(s)

    # 좌측 헤드라인
    add_rich(s, Inches(0.7), Inches(1.9), Inches(5.5), Inches(2.5),
        segs=[
            ('더 나은 금융을 위한,', 38, True, GRAY_900),
            'NL',
            ('우리', 38, True, WOORI_BLUE),
            ('의 약속', 38, True, GRAY_900),
        ], line_sp=1.25)

    add_underline(s, Inches(0.7), Inches(4.25))

    add_text(s, Inches(0.7), Inches(4.5), Inches(6), Inches(2),
        "우리는 고객 한 분 한 분의 일상에\n더 가까이 다가가, 더 편리하고 안전한 금융 경험을\n제공하기 위해 노력하고 있습니다.\n\n기술과 신뢰를 바탕으로,\n당신의 내일을 함께 만들어갑니다.",
        size=13, color=GRAY_500, line_sp=1.7)

    add_corner_decor(s)

    # 우측: "우리가 만드는 변화" + 2x2 카드
    add_text(s, Inches(7.4), Inches(1.4), Inches(5.5), Inches(0.4),
        "우리가 만드는 변화", size=14, bold=True, color=GRAY_900)
    add_line(s, Inches(9.6), Inches(1.65), Inches(3.3), Emu(8000), color=GRAY_200)

    cards = [
        ('🛡', '안전한 금융', WOORI_BLUE_LITE, WOORI_BLUE,
         '마이데이터 기반 개인화 경고와\n사전 차단으로 손실 예방',
         ['신용대출·자동이체 인용 경고', '딥페이크 송금 차단', '한도제한 + AI 보조']),
        ('📱', '편리한 금융', WOORI_BLUE_LITE, WOORI_BLUE,
         '음성과 단계별 동행으로\n진입 단계부터 풀어드립니다',
         ['음성 우선 인터페이스', '인증 단계 동행 안내', '한 번의 동의 · 끝까지']),
        ('👥', '고객 중심 금융', WOORI_BLUE_LITE, WOORI_BLUE,
         '같이 고민하는 옆자리 행원,\n결정은 항상 고객 본인',
         ['"잠깐만요" 3단계 경고', '결과 강요 X · 선택지 동등', '고령자 페르소나 토글']),
        ('🌱', '함께 성장하는 금융', WOORI_BLUE_LITE, WOORI_BLUE,
         '시니어 · 청년 · 첫 투자자까지\n모두를 위한 보호망',
         ['시니어플러스 점포 연계', '청년 첫 투자자 보호', '금감원 가이드라인 반영']),
    ]
    base_l = Inches(7.4)
    base_t = Inches(1.95)
    card_w = Inches(2.65)
    card_h = Inches(2.4)
    gap = Inches(0.18)
    for i, (icon, title, ibg, ic, desc, checks) in enumerate(cards):
        row, col = i // 2, i % 2
        x = base_l + (card_w + gap) * col
        y = base_t + (card_h + gap) * row
        add_round(s, x, y, card_w, card_h, fill=WHITE, line=GRAY_200, lw=Pt(0.5), radius=0.04)
        add_icon_box(s, x + Inches(0.2), y + Inches(0.2), Inches(0.55), ibg, ic, icon)
        add_text(s, x + Inches(0.2), y + Inches(0.88), card_w - Inches(0.4), Inches(0.35),
            title, size=15, bold=True, color=WOORI_BLUE)
        add_text(s, x + Inches(0.2), y + Inches(1.25), card_w - Inches(0.4), Inches(0.6),
            desc, size=9.5, color=GRAY_500, line_sp=1.4)
        add_dotted_line(s, x + Inches(0.2), y + Inches(1.85), card_w - Inches(0.4))
        for j, ck in enumerate(checks):
            add_check_item(s, x + Inches(0.2), y + Inches(1.95) + Inches(0.13 * j),
                card_w - Inches(0.4), ck, size=8.5, color=GRAY_500)


# ════════════════════════════════════════════════════════
# 슬라이드 9 — 우리의 약속으로 (가이드3 정확 카피 1x4)
# ════════════════════════════════════════════════════════
def slide_09():
    s = prs.slides.add_slide(blank)
    add_gradient_bg(s)
    add_header_logo(s)
    add_page_no(s, 9)
    add_footer(s)
    add_bar_chart_deco(s)

    # 좌측 큰 헤드라인 (가이드3 톤)
    add_rich(s, Inches(0.7), Inches(1.6), Inches(6), Inches(3),
        segs=[
            ('우리', 44, True, WOORI_BLUE),
            ('의 약속으로', 44, True, GRAY_900),
            'NL',
            ('더 나은 금융', 44, True, GRAY_900),
            ('을 만듭니다', 44, True, GRAY_900),
        ], line_sp=1.25)

    add_underline(s, Inches(0.7), Inches(4.4))

    add_text(s, Inches(0.7), Inches(4.7), Inches(6), Inches(1.5),
        "우리는 고객의 신뢰를 최우선으로 생각하며,\n안전하고 편리한 금융 서비스를 통해\n모두가 함께 성장하는 내일을 만들어갑니다.",
        size=13, color=GRAY_500, line_sp=1.7)

    # 좌하단 도시 + 캐릭터 자리
    add_city_silhouette(s, Inches(0), Inches(6.7), Inches(6.5), Inches(0.6))

    # 좌하단 말풍선 박스
    bubble_l, bubble_t = Inches(0.7), Inches(6.5)
    add_round(s, bubble_l, bubble_t, Inches(5.8), Inches(0.7),
        fill=WHITE, line=WOORI_BLUE_LITE, lw=Pt(1), radius=0.1)
    add_oval(s, bubble_l + Inches(0.15), bubble_t + Inches(0.15),
        Inches(0.4), Inches(0.4), WOORI_BLUE)
    add_text(s, bubble_l + Inches(0.15), bubble_t + Inches(0.15),
        Inches(0.4), Inches(0.4), "💬", size=12, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rich(s, bubble_l + Inches(0.7), bubble_t + Inches(0.13),
        Inches(5), Inches(0.5),
        segs=[
            ('옆자리 행원의 귀환  ', 12, True, WOORI_BLUE),
            ('· 우리은행은 디지털에서도 같은 자리에', 10, False, GRAY_500),
        ])
    add_text(s, bubble_l + Inches(0.7), bubble_t + Inches(0.4),
        Inches(5), Inches(0.3),
        "AI는 권유하지 않습니다. 같이 고민합니다.",
        size=10, color=GRAY_500)

    # 우측: "우리의 금융 가치는 이렇게 실현됩니다" + 1x4
    add_text(s, Inches(7.3), Inches(1.4), Inches(5.7), Inches(0.4),
        "●  WON AI 뱅커가 만들어낼 변화", size=14, bold=True, color=GRAY_900)
    add_line(s, Inches(11.0), Inches(1.65), Inches(1.9), Emu(8000), color=GRAY_200)

    items = [
        ('01', '🛡', '신뢰할 수 있는 금융', WOORI_BLUE_LITE, WOORI_BLUE,
         '본인 마이데이터를 인용해 경고\n· 출처가 분명한 사전 차단'),
        ('02', '📱', '언제 어디서나 편리한 금융', GREEN_LITE, GREEN,
         '음성·텍스트·시각 멀티 모달\n· 진입 단계부터 동행'),
        ('03', '👥', '고객 중심의 금융', PURPLE_LITE, PURPLE,
         '선택지 동등 노출 · 결과 강요 X\n· 결정은 고객 본인'),
        ('04', '🌱', '함께 성장하는 금융', WOORI_BLUE_LITE, WOORI_BLUE,
         '시니어 · 청년 · 첫 투자자\n· 모든 세대 보호망'),
    ]
    base_t = Inches(1.95)
    card_h = Inches(1.15)
    card_l = Inches(7.3)
    card_w = Inches(5.7)
    for i, (num, icon, title, ibg, ic, desc) in enumerate(items):
        y = base_t + (card_h + Inches(0.1)) * i
        add_round(s, card_l, y, card_w, card_h, fill=WHITE, line=GRAY_200, lw=Pt(0.5), radius=0.05)
        # 좌측 원형 아이콘
        add_oval(s, card_l + Inches(0.2), y + Inches(0.2), Inches(0.75), Inches(0.75), ibg)
        add_text(s, card_l + Inches(0.2), y + Inches(0.2), Inches(0.75), Inches(0.75),
            icon, size=22, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 번호 (작게)
        add_text(s, card_l + Inches(1.15), y + Inches(0.15), Inches(0.5), Inches(0.25),
            num, size=10, bold=True, color=ic)
        # 제목
        add_text(s, card_l + Inches(1.15), y + Inches(0.35), card_w - Inches(1.5), Inches(0.35),
            title, size=14, bold=True, color=GRAY_900)
        # 설명
        add_text(s, card_l + Inches(1.15), y + Inches(0.7), card_w - Inches(1.5), Inches(0.4),
            desc, size=9.5, color=GRAY_500, line_sp=1.4)
        # > 화살표
        add_text(s, card_l + card_w - Inches(0.45), y, Inches(0.35), card_h,
            "›", size=22, bold=True, color=ic,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════
# 슬라이드 10 — 기존 vs 우리 (비교표)
# ════════════════════════════════════════════════════════
def slide_10():
    s = prs.slides.add_slide(blank)
    add_gradient_bg(s)
    add_header_logo(s)
    add_page_no(s, 10)
    add_footer(s)
    add_bar_chart_deco(s)

    add_chapter_label(s, Inches(0.7), Inches(1.4), "DIFFERENCE  ·  기존 vs 우리")

    add_rich(s, Inches(0.7), Inches(1.9), Inches(12), Inches(1.2),
        segs=[
            ('다섯 영역, ', 38, True, GRAY_900),
            ('모두', 38, True, WOORI_BLUE),
            (' 다르게 합니다.', 38, True, GRAY_900),
        ])

    add_underline(s, Inches(0.7), Inches(3.0))

    add_text(s, Inches(0.7), Inches(3.2), Inches(12), Inches(0.4),
        "옆자리 행원의 역할을, 디지털 환경에 그대로 옮겼습니다.",
        size=15, color=GRAY_500)

    # 비교 표
    rows = [
        ('항목', '기존 은행앱 · 고령자 모드', 'WON AI 뱅커'),
        ('진입', '앱 설치 + 본인인증 + 로그인 통과 후', '첫 화면에서 바로 (음성 우선)'),
        ('안내', '글자만 큰 동일 화면', '음성 + 단계별 동행 · 단어 하이라이트'),
        ('경고', '"원금손실 위험" 한 줄', '본인 데이터 인용 · 3단계 사전 경고'),
        ('사기 차단', '송금 직전 차단 없음', '가족 송금 · 과소비 패턴 자동 탐지'),
        ('시니어', '별도 영업점 (오프라인)', '모바일에서도 행원 경험 · 페르소나 토글'),
    ]
    table_t = Inches(4.1)
    table_l = Inches(0.7)
    row_h = Inches(0.52)
    col_w = [Inches(2.0), Inches(4.7), Inches(5.3)]

    for ri, row in enumerate(rows):
        y = table_t + row_h * ri
        is_header = (ri == 0)
        for ci, cell in enumerate(row):
            x = table_l + sum(col_w[:ci])
            if is_header:
                add_rect(s, x, y, col_w[ci], row_h, fill=GRAY_900)
                tcolor, bold, size = WHITE, True, 13
            elif ci == 0:
                add_rect(s, x, y, col_w[ci], row_h, fill=GRAY_100,
                    line=WHITE, lw=Pt(1))
                tcolor, bold, size = GRAY_900, True, 12
            elif ci == 1:
                add_rect(s, x, y, col_w[ci], row_h, fill=WHITE,
                    line=GRAY_200, lw=Pt(0.5))
                tcolor, bold, size = GRAY_500, False, 11.5
            else:
                add_rect(s, x, y, col_w[ci], row_h, fill=WOORI_BLUE_LITE,
                    line=WOORI_BLUE_LITE, lw=Pt(0.5))
                tcolor, bold, size = WOORI_BLUE_DARK, True, 11.5
            add_text(s, x + Inches(0.2), y, col_w[ci] - Inches(0.4), row_h,
                cell, size=size, bold=bold, color=tcolor,
                anchor=MSO_ANCHOR.MIDDLE)

    # 하단 핵심 한 줄
    add_text(s, Inches(0.7), Inches(7.05), Inches(12), Inches(0.3),
        "💬  옆자리 행원의 역할을 그대로 디지털에 옮겼습니다.",
        size=13, bold=True, color=WOORI_BLUE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════
# 슬라이드 11 — 데모 ① 김영자(72)
# ════════════════════════════════════════════════════════
def slide_11():
    s = prs.slides.add_slide(blank)
    add_gradient_bg(s)
    add_header_logo(s)
    add_page_no(s, 11)
    add_footer(s)
    add_bar_chart_deco(s)

    add_chapter_label(s, Inches(0.7), Inches(1.4), "DEMO  01  ·  시니어 시나리오", color=ORANGE)

    add_rich(s, Inches(0.7), Inches(1.9), Inches(7), Inches(1.2),
        segs=[
            ('🧓  김영자 (72)', 32, True, GRAY_900),
            (' — 은퇴 교사', 18, False, GRAY_500),
        ])

    add_text(s, Inches(0.7), Inches(2.65), Inches(7), Inches(0.4),
        "노후자금 8,000만원, 첫 펀드 가입을 고민 중",
        size=15, color=GRAY_500)

    add_underline(s, Inches(0.7), Inches(3.15), color=ORANGE)

    # 페르소나 마이데이터 카드
    p_l, p_t = Inches(0.7), Inches(3.4)
    p_w, p_h = Inches(5.5), Inches(3.5)
    add_round(s, p_l, p_t, p_w, p_h, fill=WHITE, line=GRAY_200, lw=Pt(0.5), radius=0.04)
    add_text(s, p_l + Inches(0.3), p_t + Inches(0.25), Inches(4), Inches(0.4),
        "📊  마이데이터", size=13, bold=True, color=ORANGE)
    add_dotted_line(s, p_l + Inches(0.3), p_t + Inches(0.7), p_w - Inches(0.6))

    info = [
        ('나이', '72세 · 은퇴 교사'),
        ('우리은행 이용', '30년 단골'),
        ('보유 자산', '노후자금 8,000만원'),
        ('대출', '없음'),
        ('투자 경험', '없음 (첫 펀드)'),
        ('검토 펀드', '삼성 글로벌 반도체 A-e'),
        ('투자 예정', '2,000만원 (자산의 25%)'),
    ]
    for i, (k, v) in enumerate(info):
        y = p_t + Inches(0.85) + Inches(0.35 * i)
        add_text(s, p_l + Inches(0.3), y, Inches(1.7), Inches(0.3),
            k, size=11, color=GRAY_500, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, p_l + Inches(2.1), y, Inches(3.2), Inches(0.3),
            v, size=11.5, bold=True, color=GRAY_900, anchor=MSO_ANCHOR.MIDDLE)

    # 우측: 영상 자리 + 예상 결과
    v_l, v_t = Inches(6.6), Inches(1.9)
    v_w, v_h = Inches(6.2), Inches(4.0)
    add_round(s, v_l, v_t, v_w, v_h, fill=GRAY_900, line=ORANGE, lw=Pt(2), radius=0.04)
    # 재생 아이콘
    play_size = Inches(1.0)
    add_oval(s, v_l + (v_w - play_size) // 2, v_t + (v_h - play_size) // 2,
        play_size, play_size, ORANGE)
    add_text(s, v_l + (v_w - play_size) // 2, v_t + (v_h - play_size) // 2,
        play_size, play_size, "▶", size=32, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, v_l, v_t + v_h - Inches(0.45), v_w, Inches(0.35),
        "[ 김영자 시나리오 영상 — 약 60초 ]",
        size=11, color=GRAY_300, align=PP_ALIGN.CENTER)

    # 예상 결과 카드
    r_t = Inches(6.05)
    add_round(s, v_l, r_t, v_w, Inches(0.85),
        fill=ORANGE_LITE, line=ORANGE, lw=Pt(1), radius=0.05)
    add_text(s, v_l + Inches(0.3), r_t + Inches(0.1), v_w, Inches(0.3),
        "✅ 예상 결과", size=11, bold=True, color=ORANGE)
    add_text(s, v_l + Inches(0.3), r_t + Inches(0.4), v_w - Inches(0.6), Inches(0.4),
        "AI 경고 → '다시 생각해볼게요' 선택 → 원금보장형 상품으로 분기  ·  손실 0원",
        size=11, color=GRAY_700)


# ════════════════════════════════════════════════════════
# 슬라이드 12 — 데모 ② 이혜원(28)
# ════════════════════════════════════════════════════════
def slide_12():
    s = prs.slides.add_slide(blank)
    add_gradient_bg(s)
    add_header_logo(s)
    add_page_no(s, 12)
    add_footer(s)
    add_bar_chart_deco(s)

    add_chapter_label(s, Inches(0.7), Inches(1.4), "DEMO  02  ·  청년 시나리오", color=WOORI_BLUE)

    add_rich(s, Inches(0.7), Inches(1.9), Inches(7), Inches(1.2),
        segs=[
            ('👩  이혜원 (28)', 32, True, GRAY_900),
            (' — 직장 3년차', 18, False, GRAY_500),
        ])

    add_text(s, Inches(0.7), Inches(2.65), Inches(7), Inches(0.4),
        "월급 265만, 신용대출 3,000만, 카드값 6/5 결제 예정",
        size=15, color=GRAY_500)

    add_underline(s, Inches(0.7), Inches(3.15))

    # 페르소나 카드
    p_l, p_t = Inches(0.7), Inches(3.4)
    p_w, p_h = Inches(5.5), Inches(3.5)
    add_round(s, p_l, p_t, p_w, p_h, fill=WHITE, line=GRAY_200, lw=Pt(0.5), radius=0.04)
    add_text(s, p_l + Inches(0.3), p_t + Inches(0.25), Inches(4), Inches(0.4),
        "📊  마이데이터", size=13, bold=True, color=WOORI_BLUE)
    add_dotted_line(s, p_l + Inches(0.3), p_t + Inches(0.7), p_w - Inches(0.6))

    info = [
        ('나이', '28세 · 직장 3년차'),
        ('월 수입', '2,650,000원'),
        ('신용대출', '30,000,000원 (월 30만)'),
        ('카드 결제 예정', '6/5 · 423,500원'),
        ('자동이체 합계', '월 1,405,000원'),
        ('여유자금', '월 300,000원'),
        ('투자 예정', '300,000원 (여유자금 100%)'),
    ]
    for i, (k, v) in enumerate(info):
        y = p_t + Inches(0.85) + Inches(0.35 * i)
        add_text(s, p_l + Inches(0.3), y, Inches(1.9), Inches(0.3),
            k, size=11, color=GRAY_500, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, p_l + Inches(2.3), y, Inches(3.0), Inches(0.3),
            v, size=11.5, bold=True, color=GRAY_900, anchor=MSO_ANCHOR.MIDDLE)

    # 우측: 영상 + 결과
    v_l, v_t = Inches(6.6), Inches(1.9)
    v_w, v_h = Inches(6.2), Inches(4.0)
    add_round(s, v_l, v_t, v_w, v_h, fill=GRAY_900, line=WOORI_BLUE, lw=Pt(2), radius=0.04)
    play_size = Inches(1.0)
    add_oval(s, v_l + (v_w - play_size) // 2, v_t + (v_h - play_size) // 2,
        play_size, play_size, WOORI_BLUE)
    add_text(s, v_l + (v_w - play_size) // 2, v_t + (v_h - play_size) // 2,
        play_size, play_size, "▶", size=32, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, v_l, v_t + v_h - Inches(0.45), v_w, Inches(0.35),
        "[ 이혜원 시나리오 영상 — 약 60초 ]",
        size=11, color=GRAY_300, align=PP_ALIGN.CENTER)

    r_t = Inches(6.05)
    add_round(s, v_l, r_t, v_w, Inches(0.85),
        fill=WOORI_BLUE_LITE, line=WOORI_BLUE, lw=Pt(1), radius=0.05)
    add_text(s, v_l + Inches(0.3), r_t + Inches(0.1), v_w, Inches(0.3),
        "✅ 예상 결과", size=11, bold=True, color=WOORI_BLUE)
    add_text(s, v_l + Inches(0.3), r_t + Inches(0.4), v_w - Inches(0.6), Inches(0.4),
        "자동이체·카드 결제 인용 → 손실 시나리오 노출 → 본인 결정 (정보 충분히 받은 후)",
        size=11, color=GRAY_700)


# ════════════════════════════════════════════════════════
# 슬라이드 13 — 기술 아키텍처 (가이드2 패턴 카피)
# ════════════════════════════════════════════════════════
def slide_13():
    s = prs.slides.add_slide(blank)
    add_gradient_bg(s)
    add_header_logo(s)
    add_page_no(s, 13)
    add_footer(s)
    add_bar_chart_deco(s)

    add_rich(s, Inches(0.7), Inches(1.9), Inches(6), Inches(2.5),
        segs=[
            ('기술과 신뢰를', 38, True, GRAY_900),
            'NL',
            ('바탕으로,', 38, True, GRAY_900),
            'NL',
            ('우리', 38, True, WOORI_BLUE),
            ('가 만듭니다.', 38, True, GRAY_900),
        ], line_sp=1.25)

    add_underline(s, Inches(0.7), Inches(5.0))

    add_text(s, Inches(0.7), Inches(5.25), Inches(6), Inches(1.5),
        "데모는 시뮬레이션이지만,\n실제 우리은행 인프라 위에 바로 얹는 구조입니다.\n신규 인프라 투자 부담 ↓ · 기존 자산 활용 ↑",
        size=13, color=GRAY_500, line_sp=1.7)

    add_corner_decor(s)

    # 우측: 기술 카드 2x2
    add_text(s, Inches(7.4), Inches(1.4), Inches(5.5), Inches(0.4),
        "WON AI 뱅커의 기술 스택", size=14, bold=True, color=GRAY_900)
    add_line(s, Inches(10.4), Inches(1.65), Inches(2.5), Emu(8000), color=GRAY_200)

    cards = [
        ('🧠', '마이데이터 연동', WOORI_BLUE_LITE, WOORI_BLUE,
         '계좌·자동이체·카드·대출·재무목표\n실제 우리은행 인프라 위에 얹는 구조',
         ['실시간 자동 갱신', 'DSR·재무목표 자동 계산', '오픈뱅킹 확장 가능']),
        ('🎙', '음성 안내 (TTS)', GREEN_LITE, GREEN,
         'Qwen3-TTS CustomVoice "sohee"\n한국 여성 화자 + Web Speech fallback',
         ['mp3 사전 캐시 + 단어 하이라이트', '정지·재생 통제권', '발음 최적화']),
        ('⚙', '프론트엔드', PURPLE_LITE, PURPLE,
         'React 19 + Vite 8 SPA\nGitHub Pages 자동 배포',
         ['모바일 폰 프레임 390×844', '페이지 전환 0ms', '오프라인 캐시']),
        ('🛡', 'AI 경고 룰엔진', RED_LITE, RED,
         '펀드 데이터 + 사용자 마이데이터\n결합한 3단계 경고 자동 생성',
         ['DSR·과소비·재무목표 영향', '딥페이크 송금 차단 확장', 'LLM 연동 로드맵']),
    ]
    base_l = Inches(7.4)
    base_t = Inches(1.95)
    card_w = Inches(2.65)
    card_h = Inches(2.4)
    gap = Inches(0.18)
    for i, (icon, title, ibg, ic, desc, checks) in enumerate(cards):
        row, col = i // 2, i % 2
        x = base_l + (card_w + gap) * col
        y = base_t + (card_h + gap) * row
        add_round(s, x, y, card_w, card_h, fill=WHITE, line=GRAY_200, lw=Pt(0.5), radius=0.04)
        add_icon_box(s, x + Inches(0.2), y + Inches(0.2), Inches(0.55), ibg, ic, icon)
        add_text(s, x + Inches(0.2), y + Inches(0.88), card_w - Inches(0.4), Inches(0.35),
            title, size=14, bold=True, color=WOORI_BLUE)
        add_text(s, x + Inches(0.2), y + Inches(1.25), card_w - Inches(0.4), Inches(0.6),
            desc, size=9, color=GRAY_500, line_sp=1.4)
        add_dotted_line(s, x + Inches(0.2), y + Inches(1.85), card_w - Inches(0.4))
        for j, ck in enumerate(checks):
            add_check_item(s, x + Inches(0.2), y + Inches(1.95) + Inches(0.13 * j),
                card_w - Inches(0.4), ck, size=8.5, color=GRAY_500)


# ════════════════════════════════════════════════════════
# 슬라이드 14 — 사회적 가치 + 사업화 (가이드3 패턴)
# ════════════════════════════════════════════════════════
def slide_14():
    s = prs.slides.add_slide(blank)
    add_gradient_bg(s)
    add_header_logo(s)
    add_page_no(s, 14)
    add_footer(s)
    add_bar_chart_deco(s)

    add_rich(s, Inches(0.7), Inches(1.6), Inches(6), Inches(3),
        segs=[
            ('사회적 가치 + ', 36, True, GRAY_900),
            'NL',
            ('사업화 로드맵', 36, True, WOORI_BLUE),
        ], line_sp=1.25)

    add_underline(s, Inches(0.7), Inches(4.0))

    add_text(s, Inches(0.7), Inches(4.3), Inches(6), Inches(2),
        "금감원 2026 \"따뜻한 금융\" 정책과 방향성 일치.\n우리은행이 선도 모델이 되면 전 금융권\n표준 제안까지 자연스럽게 확장됩니다.",
        size=13, color=GRAY_500, line_sp=1.7)

    add_city_silhouette(s, Inches(0), Inches(6.7), Inches(6.5), Inches(0.6))

    # 좌하단 말풍선 (가이드3 톤)
    bubble_l, bubble_t = Inches(0.7), Inches(6.4)
    add_round(s, bubble_l, bubble_t, Inches(5.8), Inches(0.8),
        fill=WHITE, line=WOORI_BLUE_LITE, lw=Pt(1), radius=0.1)
    add_oval(s, bubble_l + Inches(0.15), bubble_t + Inches(0.18),
        Inches(0.4), Inches(0.4), WOORI_BLUE)
    add_text(s, bubble_l + Inches(0.15), bubble_t + Inches(0.18),
        Inches(0.4), Inches(0.4), "💬", size=12, color=WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, bubble_l + Inches(0.7), bubble_t + Inches(0.13),
        Inches(5), Inches(0.3),
        "정부 정책과 정합", size=12, bold=True, color=WOORI_BLUE)
    add_text(s, bubble_l + Inches(0.7), bubble_t + Inches(0.4),
        Inches(5), Inches(0.4),
        "다크패턴 방지 가이드라인(2026.04) · 사망자 명의도용 차단 일 1회",
        size=10, color=GRAY_500)

    # 우측 1x4 로드맵
    add_text(s, Inches(7.3), Inches(1.4), Inches(5.7), Inches(0.4),
        "●  4단계 사업화 로드맵", size=14, bold=True, color=GRAY_900)
    add_line(s, Inches(10.0), Inches(1.65), Inches(2.9), Emu(8000), color=GRAY_200)

    stages = [
        ('NOW', '🚀', 'WON AI 뱅커 v1', WOORI_BLUE_LITE, WOORI_BLUE,
         '펀드 가입 플로우 · 이혜원·김영자 페르소나 토글 시연'),
        ('+3M', '🛡', '시니어 사기 차단 확장', RED_LITE, RED,
         '딥페이크 송금 직전 차단 · 가족 위임 인증 도입'),
        ('+6M', '🏢', '시니어플러스 통합', GREEN_LITE, GREEN,
         '물리 행원 + AI 행원 동일 컨텍스트 공유 · 옴니채널'),
        ('+12M', '🌐', '전 금융권 표준 제안', PURPLE_LITE, PURPLE,
         '금감원 가이드라인 반영 모델 · 우리은행 선도'),
    ]
    base_t = Inches(1.95)
    card_h = Inches(1.18)
    card_l = Inches(7.3)
    card_w = Inches(5.7)
    for i, (stage, icon, title, ibg, ic, desc) in enumerate(stages):
        y = base_t + (card_h + Inches(0.08)) * i
        add_round(s, card_l, y, card_w, card_h, fill=WHITE, line=GRAY_200, lw=Pt(0.5), radius=0.05)
        # 좌측 원형 아이콘
        add_oval(s, card_l + Inches(0.2), y + Inches(0.2), Inches(0.78), Inches(0.78), ibg)
        add_text(s, card_l + Inches(0.2), y + Inches(0.2), Inches(0.78), Inches(0.78),
            icon, size=22, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 스테이지 배지
        add_round(s, card_l + Inches(1.15), y + Inches(0.18), Inches(0.7), Inches(0.3),
            fill=ic, radius=0.4)
        add_text(s, card_l + Inches(1.15), y + Inches(0.18), Inches(0.7), Inches(0.3),
            stage, size=9, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 제목
        add_text(s, card_l + Inches(1.95), y + Inches(0.15), card_w - Inches(2.2), Inches(0.4),
            title, size=15, bold=True, color=GRAY_900)
        # 설명
        add_text(s, card_l + Inches(1.15), y + Inches(0.6), card_w - Inches(1.4), Inches(0.55),
            desc, size=10, color=GRAY_500, line_sp=1.4)
        # > 화살표
        add_text(s, card_l + card_w - Inches(0.45), y, Inches(0.35), card_h,
            "›", size=22, bold=True, color=ic,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ════════════════════════════════════════════════════════
# 슬라이드 15 — 마무리
# ════════════════════════════════════════════════════════
def slide_15():
    s = prs.slides.add_slide(blank)
    add_gradient_bg(s)
    add_header_logo(s)
    add_page_no(s, 15)
    add_footer(s)
    add_bar_chart_deco(s)

    # 중앙 큰 카피
    add_text(s, Inches(0.7), Inches(2.5), Inches(12), Inches(1.0),
        "권유하지 않습니다.",
        size=64, bold=True, color=GRAY_900, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(3.5), Inches(12), Inches(1.0),
        "같이 고민합니다.",
        size=64, bold=True, color=WOORI_BLUE, align=PP_ALIGN.CENTER)

    # 강조선 (중앙)
    add_rect(s, Inches(6.4), Inches(4.8), Inches(0.5), Emu(40000), fill=WOORI_BLUE)

    # 부연
    add_text(s, Inches(0.7), Inches(5.1), Inches(12), Inches(0.5),
        "옆자리 행원이 사라진 자리에, 우리가 다시 앉습니다.",
        size=20, color=GRAY_500, align=PP_ALIGN.CENTER)

    # Q&A
    add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.5),
        "Thank You",
        size=18, bold=True, color=WOORI_BLUE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.4),
        "WON AI 뱅커  ·  팀 [팀명]  ·  Q & A",
        size=12, color=GRAY_500, align=PP_ALIGN.CENTER)


# ── 모든 슬라이드 ───────────────────────────────────────
makers = [slide_01, slide_02, slide_03, slide_04, slide_05,
          slide_06, slide_07, slide_08, slide_09, slide_10,
          slide_11, slide_12, slide_13, slide_14, slide_15]
for i, m in enumerate(makers, 1):
    print(f"[{i:2d}/15] {m.__name__}")
    m()

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'WON_AI_Banker.pptx')
prs.save(OUT)
print(f"\n[OK] {OUT}  ·  총 {len(prs.slides)}장")
